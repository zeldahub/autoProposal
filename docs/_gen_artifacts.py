"""Generate additional artifacts:
- screen-design.pptx
- db-design.xlsx
- interface-spec.xlsx
- api-architecture.xlsx

Run:  python _gen_artifacts.py
"""
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt

DOCS_DIR = Path(__file__).parent

# ── 공통 색상 ──────────────────────────────────────────────
NAVY = "1F3A5F"
NAVY_RGB = PPTRGB(0x1F, 0x3A, 0x5F)
GRAY = "F4F4F4"
WHITE = "FFFFFF"
ACCENT = "3B82F6"
ACCENT_RGB = PPTRGB(0x3B, 0x82, 0xF6)
ROW_ALT = "F7FAFC"

KO_FONT = "맑은 고딕"

# ── XLSX 헬퍼 ─────────────────────────────────────────────
THIN = Side(style="thin", color="BBBBBB")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


def style_header(cell):
    cell.font = Font(name=KO_FONT, bold=True, color=WHITE, size=11)
    cell.fill = PatternFill("solid", fgColor=NAVY)
    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    cell.border = BORDER


def style_body(cell, alt=False):
    cell.font = Font(name=KO_FONT, size=10)
    cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
    cell.border = BORDER
    if alt:
        cell.fill = PatternFill("solid", fgColor=ROW_ALT)


def write_table(ws, start_row, headers, rows, widths=None):
    for c, h in enumerate(headers, start=1):
        cell = ws.cell(row=start_row, column=c, value=h)
        style_header(cell)
    for r, row in enumerate(rows, start=start_row + 1):
        alt = ((r - start_row) % 2 == 0)
        for c, v in enumerate(row, start=1):
            cell = ws.cell(row=r, column=c, value=v)
            style_body(cell, alt=alt)
    if widths:
        for i, w in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w
    ws.row_dimensions[start_row].height = 24
    ws.freeze_panes = ws.cell(row=start_row + 1, column=1)
    return start_row + 1 + len(rows)


def add_title(ws, text, sub=None):
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=8)
    c = ws.cell(row=1, column=1, value=text)
    c.font = Font(name=KO_FONT, bold=True, size=18, color=NAVY)
    c.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 30
    if sub:
        ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=8)
        c2 = ws.cell(row=2, column=1, value=sub)
        c2.font = Font(name=KO_FONT, size=10, color="666666")
        c2.alignment = Alignment(horizontal="left", vertical="center")
        return 4
    return 3


# ╔══════════════════════════════════════════════════════════
# ║ 1) screen-design.pptx
# ╚══════════════════════════════════════════════════════════
def gen_screen_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def title_box(slide, text, subtitle=None):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = KO_FONT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY_RGB
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = KO_FONT
            p2.font.size = Pt(13)
            p2.font.color.rgb = PPTRGB(0x66, 0x66, 0x66)
        # underline bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2.0), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_RGB
        bar.line.fill.background()

    def add_table_pptx(slide, left, top, width, height, headers, rows):
        cols = len(headers)
        tbl_shape = slide.shapes.add_table(rows=len(rows) + 1, cols=cols, left=left, top=top, width=width, height=height)
        tbl = tbl_shape.table
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY_RGB
            tf = cell.text_frame
            tf.text = h
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.bold = True
                    r.font.color.rgb = PPTRGB(0xFF, 0xFF, 0xFF)
                    r.font.name = KO_FONT
                    r.font.size = Pt(11)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                if r % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = PPTRGB(0xF7, 0xFA, 0xFC)
                tf = cell.text_frame
                tf.text = str(val) if val is not None else ""
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.name = KO_FONT
                        run.font.size = Pt(10)
        return tbl

    def add_text_block(slide, left, top, width, height, lines, font_size=12, mono=False):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = line
            p.font.name = "Consolas" if mono else KO_FONT
            p.font.size = Pt(font_size)
            p.font.color.rgb = PPTRGB(0x33, 0x33, 0x33)

    def footer(slide, page, total):
        box = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Lon · 화면설계서  |  {page}/{total}"
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(9)
        p.font.name = KO_FONT
        p.font.color.rgb = PPTRGB(0x99, 0x99, 0x99)

    # ── 표지 ──
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PPTRGB(0x0F, 0x16, 0x23)
    bg.line.fill.background()
    title = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Lon — AI 사업제안서 자동 생성기"
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = PPTRGB(0xFF, 0xFF, 0xFF); p.font.name = KO_FONT
    p2 = tf.add_paragraph(); p2.text = "화면 설계서 (Screen Design)"
    p2.font.size = Pt(24); p2.font.color.rgb = PPTRGB(0x3B, 0x82, 0xF6); p2.font.name = KO_FONT
    meta = s.shapes.add_textbox(Inches(1), Inches(6.0), Inches(11.3), Inches(0.7))
    pm = meta.text_frame.paragraphs[0]
    pm.text = "v0.1 (초안)   |   2026-05-03"
    pm.font.size = Pt(13); pm.font.color.rgb = PPTRGB(0xCC, 0xCC, 0xCC); pm.font.name = KO_FONT

    slides_total = 10  # placeholder; updated later

    # ── 화면 목록 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "1. 화면 목록", "전체 11화면, 권한 분리(USER/ADMIN)")
    add_table_pptx(
        s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.2),
        headers=["ID", "화면명", "URL", "설명", "권한"],
        rows=[
            ["S-000", "로그인", "/login", "로컬 계정 로그인", "비회원"],
            ["S-010", "홈/대시보드", "/", "최근 사업, 산출물 통계", "사용자"],
            ["S-100", "사업제안서 생성기(메인)", "/generator", "이미지 기준 메인 화면", "사용자"],
            ["S-110", "사업 목록", "/projects", "저장된 사업 목록/검색", "사용자"],
            ["S-111", "사업 상세", "/projects/:id", "사업 정보 + 산출물 이력", "사용자"],
            ["S-120", "산출물 미리보기", "/projects/:id/preview", "PPTX/XLSX 미리보기", "사용자"],
            ["S-200", "AI 키 관리", "/settings/ai", "Provider별 키 등록·검증", "사용자"],
            ["S-210", "환경 설정", "/settings/env", "언어/테마/저장 경로", "사용자"],
            ["S-300", "표준 목차 관리", "/admin/category", "제안서 카테고리 마스터", "관리자"],
            ["S-310", "사용자 관리", "/admin/users", "계정/권한 관리", "관리자"],
            ["S-900", "에러/접근 거부", "/error/:code", "4xx/5xx 공통 화면", "공통"],
        ],
    )
    footer(s, 2, slides_total)

    # ── 공통 영역 / 디자인 토큰 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "2. 공통 영역 & 디자인 토큰", "다크 테마 / 카드형 UI")
    add_text_block(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.0), [
        "[Top Bar]",
        "  로고 · 사업명 · 사용자메뉴 · 알림 · 테마",
        "",
        "[Side]                 [Content]",
        "  생성기                 ① AI 선택",
        "  사업                   ② 분석",
        "  설정                   ③ 정보 입력",
        "                        ④ 산출",
        "                        ⑤ 카테고리",
        "",
        "[Footer] © Lon 2026",
    ], font_size=12, mono=True)
    add_table_pptx(
        s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(4.5),
        headers=["토큰", "값", "용도"],
        rows=[
            ["--color-bg", "#0F1623", "다크 배경"],
            ["--color-surface", "#1A2332", "카드 배경"],
            ["--color-primary", "#3B82F6", "CTA 버튼"],
            ["--color-accent", "#10B981", "성공/생성"],
            ["--color-danger", "#EF4444", "에러/Not Found"],
            ["--radius-md", "12px", "카드 라운드"],
            ["--font-ko", "맑은 고딕", "본문"],
        ],
    )
    footer(s, 3, slides_total)

    # ── S-100 메인 와이어프레임 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "3. S-100 사업제안서 생성기 — 와이어프레임", "이미지 5개 영역 매핑 (① ~ ⑤)")
    # 좌측: 와이어 박스들
    def section_box(left, top, width, height, label, color=PPTRGB(0x3B, 0x82, 0xF6)):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = PPTRGB(0x1A, 0x23, 0x32)
        sh.line.color.rgb = color; sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = KO_FONT
        p.font.size = Pt(11)
        p.font.color.rgb = PPTRGB(0xFF, 0xFF, 0xFF)
        return sh

    section_box(Inches(0.5), Inches(1.4), Inches(7.5), Inches(0.8), "① AI 서비스 선택  [ChatGPT][Gemini✓][Claude]  Key/Model  [테스트][상태]")
    section_box(Inches(0.5), Inches(2.3), Inches(7.5), Inches(0.9), "② 데이터 파일(MD) 분석  [공고문 Drop] [관련산출물 Drop]  [분석 시작]")
    section_box(Inches(0.5), Inches(3.3), Inches(7.5), Inches(2.0), "③ 사업 정보 입력  (회사명/사업명/목표/범위/일정/조직/인력/비용 등 12 필드)")
    section_box(Inches(0.5), Inches(5.4), Inches(7.5), Inches(0.7), "④ 산출 사항  ☑ PPTX  ☑ XLSX   [PPTX 생성] [WBS 생성]")
    section_box(Inches(0.5), Inches(6.2), Inches(7.5), Inches(0.7), "⑤ 표준 목차 카테고리 카드 그리드 (사업개요·일반사항·기술요구·…)")

    # 우측: 컴포넌트 매핑 표
    add_table_pptx(
        s, Inches(8.2), Inches(1.4), Inches(4.6), Inches(5.5),
        headers=["영역", "컴포넌트", "이벤트"],
        rows=[
            ["①", "ProviderCard", "onChange→리셋"],
            ["①", "ApiKeyInput", "onBlur→마스킹"],
            ["①", "BtnTest", "→ /api/llm/test"],
            ["②", "Dropzone×2", "onDrop→검증"],
            ["②", "BtnAnalyze", "→ /files/analyze"],
            ["③", "ProjectForm", "onChange→자동저장"],
            ["④", "BtnGenPPTX", "→ /generate/pptx"],
            ["④", "BtnGenWBS", "→ /generate/wbs"],
            ["⑤", "CategoryCard", "→ 상세 모달"],
        ],
    )
    footer(s, 4, slides_total)

    # ── 컴포넌트 명세 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "4. S-100 컴포넌트 상세 명세")
    add_table_pptx(
        s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
        headers=["코드", "컴포넌트", "타입", "필수", "검증", "이벤트"],
        rows=[
            ["C-101", "ProviderCard", "라디오(카드3)", "Y", "1개 선택", "onChange"],
            ["C-102", "ApiKeyInput", "password", "Y", "길이≥10", "onBlur"],
            ["C-103", "ModelSelect", "select", "Y", "provider별 옵션", "onChange"],
            ["C-104", "BtnTest", "button", "-", "키 입력시 활성", "→/api/llm/test"],
            ["C-201", "UploadDropzone×2", "drag&drop", "N", "≤10MB, ext 화이트리스트", "onDrop"],
            ["C-202", "BtnAnalyze", "button", "-", "첨부 1개 이상", "→/api/files/analyze"],
            ["C-301", "TextField×7", "input", "일부", "길이/형식", "자동저장"],
            ["C-302", "TextArea×7", "textarea", "일부", "1000자", "자동저장"],
            ["C-401", "Checkbox×2", "checkbox", "-", "1개 이상", "onChange"],
            ["C-402", "BtnGenPPTX", "button", "-", "④ 1개 체크", "→/api/generate/pptx"],
            ["C-403", "BtnGenWBS", "button", "-", "④ 1개 체크", "→/api/generate/wbs"],
            ["C-501", "CategoryCard×N", "card", "-", "-", "→상세 모달"],
        ],
    )
    footer(s, 5, slides_total)

    # ── 검증 규칙 / 상태 머신 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "5. 검증 규칙 & 상태 머신")
    add_table_pptx(
        s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.5),
        headers=["대상", "규칙"],
        rows=[
            ["회사명/사업명", "1~80자, 특수문자 < > 금지"],
            ["API Key", "provider별 정규식 (sk-, AIza, sk-ant-)"],
            ["첨부 파일", "PDF/DOCX/TXT/MD only, ≤10MB / ≤30MB(합)"],
            ["추진 예산", "숫자 + 단위 (예: 1,200,000,000원)"],
            ["TextArea", "1000자, XSS 이스케이프"],
        ],
    )
    add_text_block(s, Inches(8.3), Inches(1.4), Inches(4.7), Inches(4.5), [
        "[생성 버튼 상태 머신]",
        "",
        "  idle → click",
        "    → validating",
        "       → OK → generating",
        "         → done → success → download",
        "         → fail → error",
        "       → fail → error",
    ], font_size=12, mono=True)
    footer(s, 6, slides_total)

    # ── 사용자 흐름 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "6. 사용자 흐름 (User Flow)")
    flow = [
        ("1. 로그인", PPTRGB(0x10, 0xB9, 0x81)),
        ("2. 키 검증", PPTRGB(0x3B, 0x82, 0xF6)),
        ("3. 첨부 분석", PPTRGB(0x3B, 0x82, 0xF6)),
        ("4. 폼 보완", PPTRGB(0x3B, 0x82, 0xF6)),
        ("5. PPTX 생성", PPTRGB(0xEF, 0x44, 0x44)),
        ("6. 미리보기/다운로드", PPTRGB(0x10, 0xB9, 0x81)),
    ]
    x = Inches(0.4)
    for i, (label, color) in enumerate(flow):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(1.95), Inches(1.0))
        box.fill.solid(); box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = label
        p.font.name = KO_FONT; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = PPTRGB(0xFF,0xFF,0xFF)
        x += Inches(2.1)
        if i < len(flow) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x - Inches(0.15), Inches(2.85), Inches(0.15), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = PPTRGB(0x99,0x99,0x99)
            arr.line.fill.background()
    add_text_block(s, Inches(0.5), Inches(4.3), Inches(12), Inches(2),
                   ["기존 사업 재활용:  S-110 검색 → S-111 상세 → [복제] → S-100 사전채움 → ④ 재생성"],
                   font_size=14)
    footer(s, 7, slides_total)

    # ── 권한별 메뉴 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "7. 권한별 메뉴 노출")
    add_table_pptx(
        s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5),
        headers=["메뉴", "USER", "ADMIN", "비로그인"],
        rows=[
            ["1. 홈", "✓", "✓", "–"],
            ["2. 사업제안서 생성", "✓", "✓", "–"],
            ["3. 사업 관리", "✓", "✓", "–"],
            ["4. 산출물 라이브러리", "✓", "✓", "–"],
            ["5. 설정", "✓", "✓", "–"],
            ["6. 관리", "–", "✓", "–"],
            ["7. 도움말", "✓", "✓", "✓"],
            ["로그인", "–", "–", "✓"],
        ],
    )
    footer(s, 8, slides_total)

    # ── 반응형/접근성 / 에러 표시 ──
    s = prs.slides.add_slide(blank)
    title_box(s, "8. 반응형·접근성 / 에러 표시")
    add_table_pptx(
        s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(4.0),
        headers=["분류", "정책"],
        rows=[
            ["Breakpoint", "≥1280 3-col / 768~1279 2-col / <768 1-col"],
            ["키보드", "Tab 순서 보장, ESC 모달 닫기"],
            ["명도 대비", "WCAG AA, 색단독 의미금지(아이콘 병기)"],
            ["폰트 fallback", "맑은 고딕 → Malgun Gothic → Apple SD"],
        ],
    )
    add_table_pptx(
        s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(4.0),
        headers=["케이스", "UI"],
        rows=[
            ["API 키 무효", "Badge red + 토스트 + 입력 적색 보더"],
            ["파일 초과", "Dropzone 내부 inline"],
            ["LLM 타임아웃", "모달 + 재시도"],
            ["서버 5xx", "S-900 + traceId"],
        ],
    )
    footer(s, 9, slides_total)

    # ── 마무리 ──
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PPTRGB(0x0F, 0x16, 0x23); bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]; p.text = "End of Screen Design"; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = PPTRGB(0xFF,0xFF,0xFF); p.font.name = KO_FONT
    p2 = tf.add_paragraph(); p2.text = "Next: DB 설계서 (db-design.xlsx) / 연계 정의서 (interface-spec.xlsx)"
    p2.alignment = PP_ALIGN.CENTER; p2.font.size = Pt(14); p2.font.color.rgb = PPTRGB(0xAA,0xAA,0xAA); p2.font.name = KO_FONT

    out = DOCS_DIR / "screen-design.pptx"
    prs.save(out)
    print(f"OK  -> {out.name}  ({len(prs.slides)} slides)")


# ╔══════════════════════════════════════════════════════════
# ║ 2) db-design.xlsx
# ╚══════════════════════════════════════════════════════════
def gen_db_xlsx():
    wb = Workbook()
    # 표지/요약
    ws = wb.active
    ws.title = "00_요약"
    r = add_title(ws, "Lon — DB 설계서 (Hybrid)", "MariaDB 11.4 + MongoDB 8.2 / 작성일 2026-05-03")
    r = write_table(ws, r, ["구분", "내용"], [
        ["설계 원칙", "정형(메타) → MariaDB / 비정형(LLM 본문, 추출 텍스트) → MongoDB"],
        ["문자셋", "MariaDB utf8mb4_unicode_ci / Mongo UTF-8"],
        ["시간", "UTC 저장, 표시 시 KST 변환"],
        ["식별자", "MariaDB BIGINT AI + uuid CHAR(36) / Mongo ObjectId + uuid 인덱스"],
        ["논리삭제", "deleted_at TIMESTAMP / deletedAt ISODate"],
        ["일관성", "Outbox + 보상잡(mongo_repair) — 분산 트랜잭션 미사용"],
    ], widths=[18, 90])

    # MariaDB 테이블 일람
    ws = wb.create_sheet("01_MariaDB_일람")
    r = add_title(ws, "MariaDB 테이블 일람", "총 8개")
    r = write_table(ws, r, ["#", "테이블", "용도", "주요 인덱스"], [
        [1, "user", "사용자/권한", "uq email"],
        [2, "project", "사업 메타", "owner_id, status, FT(name)"],
        [3, "project_attachment", "첨부 파일 메타", "project_id, uq sha256"],
        [4, "artifact", "산출물 메타(PPTX/XLSX)", "uq (project,type,version)"],
        [5, "ai_provider_setting", "사용자별 AI 키/모델", "uq (user,provider,alias)"],
        [6, "llm_call_log", "LLM 호출 메타(본문은 Mongo)", "project_id, purpose"],
        [7, "proposal_category", "표준 목차 마스터", "uq code"],
        [8, "audit_log", "감사 로그", "user_id, action"],
    ], widths=[5, 26, 36, 30])

    # 테이블별 컬럼 명세
    columns = {
        "T01_user": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["uuid", "CHAR(36)", "Y", "UUID()", "외부 키"],
            ["email", "VARCHAR(255)", "Y", "-", "로그인 ID, UNIQUE"],
            ["password_hash", "VARCHAR(255)", "Y", "-", "bcrypt"],
            ["display_name", "VARCHAR(100)", "N", "-", "표시명"],
            ["role", "ENUM(USER,ADMIN)", "Y", "USER", "권한"],
            ["last_login_at", "TIMESTAMP", "N", "-", ""],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
            ["updated_at", "TIMESTAMP", "Y", "ON UPDATE", ""],
            ["deleted_at", "TIMESTAMP", "N", "NULL", "논리삭제"],
        ],
        "T02_project": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["uuid", "CHAR(36)", "Y", "UUID()", "UNIQUE — Mongo 외부 키"],
            ["owner_id", "BIGINT", "Y", "-", "FK→user.id"],
            ["company_name", "VARCHAR(120)", "N", "-", ""],
            ["project_name", "VARCHAR(200)", "Y", "-", "사업명, FT(ngram)"],
            ["goal", "TEXT", "N", "-", "사업 목표"],
            ["scope", "TEXT", "N", "-", "사업 범위"],
            ["schedule", "TEXT", "N", "-", "일정"],
            ["organization", "TEXT", "N", "-", "수행 조직"],
            ["staff", "TEXT", "N", "-", "수행 인력"],
            ["cost_dev", "TEXT", "N", "-", "개발 비용"],
            ["cost_ops", "TEXT", "N", "-", "운영 비용"],
            ["license_info", "TEXT", "N", "-", "라이선스"],
            ["availability", "TEXT", "N", "-", "가용성"],
            ["budget", "VARCHAR(50)", "N", "-", "예산"],
            ["ai_provider", "ENUM", "N", "-", "OPENAI/GEMINI/ANTHROPIC"],
            ["ai_model", "VARCHAR(80)", "N", "-", ""],
            ["status", "ENUM", "Y", "DRAFT", "DRAFT/READY/GENERATED/ARCHIVED"],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
            ["updated_at", "TIMESTAMP", "Y", "ON UPDATE", ""],
            ["deleted_at", "TIMESTAMP", "N", "NULL", "논리삭제"],
        ],
        "T03_project_attachment": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["project_id", "BIGINT", "Y", "-", "FK"],
            ["slot", "ENUM(NOTICE,REFERENCE)", "Y", "-", "사업공고/관련 산출물"],
            ["filename", "VARCHAR(255)", "Y", "-", ""],
            ["mime_type", "VARCHAR(100)", "Y", "-", ""],
            ["size_bytes", "INT", "Y", "-", ""],
            ["sha256", "CHAR(64)", "Y", "-", "중복/무결성"],
            ["storage_path", "VARCHAR(500)", "Y", "-", "로컬 경로"],
            ["mongo_doc_id", "CHAR(24)", "N", "-", "Mongo documents._id"],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
        ],
        "T04_artifact": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["project_id", "BIGINT", "Y", "-", "FK"],
            ["type", "ENUM(PPTX,XLSX)", "Y", "-", ""],
            ["version", "INT", "Y", "1", "1부터 증가"],
            ["filename", "VARCHAR(255)", "Y", "-", ""],
            ["storage_path", "VARCHAR(500)", "Y", "-", ""],
            ["size_bytes", "INT", "Y", "-", ""],
            ["sha256", "CHAR(64)", "Y", "-", ""],
            ["llm_call_log_id", "BIGINT", "N", "-", "어느 호출이 만들었는지"],
            ["mongo_draft_id", "CHAR(24)", "N", "-", "Mongo proposalDrafts._id"],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
        ],
        "T05_ai_provider_setting": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["user_id", "BIGINT", "Y", "-", "FK"],
            ["provider", "ENUM", "Y", "-", "OPENAI/GEMINI/ANTHROPIC"],
            ["alias", "VARCHAR(80)", "N", "-", "별칭"],
            ["api_key_cipher", "VARBINARY(512)", "Y", "-", "AES-256-GCM 암호문"],
            ["key_iv", "VARBINARY(16)", "Y", "-", "IV"],
            ["key_tag", "VARBINARY(16)", "Y", "-", "GCM tag"],
            ["default_model", "VARCHAR(80)", "N", "-", ""],
            ["temperature", "DECIMAL(3,2)", "N", "0.40", ""],
            ["max_tokens", "INT", "N", "-", ""],
            ["is_active", "TINYINT(1)", "Y", "1", ""],
            ["last_verified_at", "TIMESTAMP", "N", "-", ""],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
            ["updated_at", "TIMESTAMP", "Y", "ON UPDATE", ""],
        ],
        "T06_llm_call_log": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["project_id", "BIGINT", "Y", "-", "FK"],
            ["provider", "ENUM", "Y", "-", ""],
            ["model", "VARCHAR(80)", "Y", "-", ""],
            ["purpose", "ENUM", "Y", "-", "ANALYZE/GEN_PPTX/GEN_WBS/TEST"],
            ["input_tokens", "INT", "N", "-", ""],
            ["output_tokens", "INT", "N", "-", ""],
            ["latency_ms", "INT", "N", "-", ""],
            ["http_status", "SMALLINT", "N", "-", ""],
            ["error_code", "VARCHAR(50)", "N", "-", ""],
            ["mongo_session_id", "CHAR(24)", "N", "-", "Mongo llmSessions._id"],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
        ],
        "T07_proposal_category": [
            ["id", "INT AI", "Y", "-", "PK"],
            ["code", "VARCHAR(40)", "Y", "-", "UNIQUE (예: OVERVIEW)"],
            ["name_ko", "VARCHAR(80)", "Y", "-", "사업 개요"],
            ["parent_id", "INT", "N", "-", "자기참조"],
            ["sort_order", "INT", "Y", "-", ""],
            ["slide_template_key", "VARCHAR(80)", "N", "-", "python-pptx 템플릿 키"],
            ["system_prompt", "TEXT", "N", "-", "LLM 시스템 프롬프트"],
            ["is_active", "TINYINT(1)", "Y", "1", ""],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
            ["updated_at", "TIMESTAMP", "Y", "ON UPDATE", ""],
        ],
        "T08_audit_log": [
            ["id", "BIGINT AI", "Y", "-", "PK"],
            ["user_id", "BIGINT", "N", "-", ""],
            ["action", "VARCHAR(80)", "Y", "-", "예: PROJECT.CREATE"],
            ["target_type", "VARCHAR(40)", "N", "-", ""],
            ["target_uuid", "CHAR(36)", "N", "-", ""],
            ["ip", "VARCHAR(45)", "N", "-", ""],
            ["user_agent", "VARCHAR(255)", "N", "-", ""],
            ["meta_json", "JSON", "N", "-", "부가 정보"],
            ["created_at", "TIMESTAMP", "Y", "CURRENT_TIMESTAMP", ""],
        ],
    }
    for sheet, rows in columns.items():
        ws = wb.create_sheet(sheet)
        r = add_title(ws, f"MariaDB · {sheet[4:]}", "컬럼 명세")
        write_table(ws, r, ["컬럼", "타입", "NN", "기본", "설명"], rows,
                    widths=[24, 24, 6, 22, 50])

    # MongoDB 컬렉션 일람
    ws = wb.create_sheet("10_Mongo_일람")
    r = add_title(ws, "MongoDB 컬렉션 일람", "총 6개 / DB: lon")
    write_table(ws, r, ["#", "컬렉션", "용도", "주요 인덱스"], [
        [1, "documents", "추출 텍스트, 청크, 임베딩(옵션)", "{projectUuid:1}, {attachmentId:1}"],
        [2, "analysisResults", "공고문 분석 → 사업정보 자동 채움", "{projectUuid:1}"],
        [3, "llmSessions", "LLM 호출 요청/응답 원문", "{projectUuid:1, createdAt:-1}, TTL 90d"],
        [4, "proposalDrafts", "슬라이드별 편집 가능 초안", "{projectUuid:1, version:-1}"],
        [5, "wbsTasks", "WBS 작업 트리", "{projectUuid:1, version:-1}"],
        [6, "categoryPrompts", "카테고리별 시스템 프롬프트 이력", "{code:1, version:-1}"],
    ], widths=[5, 22, 40, 38])

    # 컬렉션 스키마 (요약)
    mongo_schemas = {
        "M01_documents": [
            ["_id", "ObjectId", "PK"],
            ["projectUuid", "string", "MariaDB project.uuid"],
            ["attachmentId", "long", "MariaDB project_attachment.id"],
            ["slot", "enum", "NOTICE | REFERENCE"],
            ["filename", "string", ""],
            ["mimeType", "string", ""],
            ["extractedText", "string", "전체 추출 텍스트"],
            ["chunks[]", "array<object>", "{idx, text, tokens, embedding?}"],
            ["summary", "string", "요약 1~2단락"],
            ["language", "string", "ko/en"],
            ["createdAt", "date", ""],
        ],
        "M02_analysisResults": [
            ["_id", "ObjectId", "PK"],
            ["projectUuid", "string", ""],
            ["source.noticeDocId", "ObjectId", ""],
            ["source.referenceDocIds", "array<ObjectId>", ""],
            ["fields.projectName", "string", "에코드림(ecoDream)"],
            ["fields.goal/scope/...", "string", "12 필드"],
            ["confidence.{field}", "double", "신뢰도 (0~1)"],
            ["model", "string", "gemini-2.5-flash"],
            ["createdAt", "date", ""],
        ],
        "M03_llmSessions": [
            ["_id", "ObjectId", "PK"],
            ["projectUuid", "string", ""],
            ["purpose", "enum", "ANALYZE/GEN_PPTX/GEN_WBS/TEST"],
            ["provider", "enum", "OPENAI/GEMINI/ANTHROPIC"],
            ["model", "string", ""],
            ["request.system", "string", ""],
            ["request.messages", "array", ""],
            ["request.temperature", "double", ""],
            ["response.text", "string", ""],
            ["response.raw", "object", "provider raw json"],
            ["usage.input/output", "int", "토큰"],
            ["latencyMs", "int", ""],
            ["createdAt", "date", "TTL 90d"],
        ],
        "M04_proposalDrafts": [
            ["_id", "ObjectId", "PK"],
            ["projectUuid", "string", ""],
            ["version", "int", "1부터"],
            ["categories[].code", "string", "OVERVIEW 등"],
            ["categories[].name", "string", ""],
            ["categories[].slides[]", "array<object>", "{title,bullets,speakerNote}"],
            ["model", "string", ""],
            ["createdAt", "date", ""],
        ],
        "M05_wbsTasks": [
            ["_id", "ObjectId", "PK"],
            ["projectUuid", "string", ""],
            ["version", "int", ""],
            ["phases[].code", "string", "P1 등"],
            ["phases[].name", "string", ""],
            ["phases[].tasks[]", "array<object>", "{code,name,owner,durationDays,deliverables}"],
            ["totalTasks", "int", ""],
            ["createdAt", "date", ""],
        ],
        "M06_categoryPrompts": [
            ["_id", "ObjectId", "PK"],
            ["code", "string", "OVERVIEW 등"],
            ["version", "int", "버전"],
            ["systemPrompt", "string", ""],
            ["userPromptTemplate", "string", ""],
            ["active", "bool", "현재 활성"],
            ["createdAt", "date", ""],
        ],
    }
    for sheet, rows in mongo_schemas.items():
        ws = wb.create_sheet(sheet)
        r = add_title(ws, f"MongoDB · {sheet[4:]}", "필드 명세")
        write_table(ws, r, ["필드", "타입", "설명"], rows, widths=[28, 22, 60])

    # 인덱스 요약
    ws = wb.create_sheet("20_인덱스")
    r = add_title(ws, "인덱스 요약", "MariaDB + MongoDB")
    write_table(ws, r, ["저장소", "대상", "인덱스", "비고"], [
        ["MariaDB", "user", "uq(email)", ""],
        ["MariaDB", "project", "idx(owner_id), idx(status), FT(project_name) WITH PARSER ngram", "한글 검색"],
        ["MariaDB", "project_attachment", "idx(project_id), uq(project_id, sha256)", "중복 방지"],
        ["MariaDB", "artifact", "uq(project_id, type, version)", "버전관리"],
        ["MariaDB", "ai_provider_setting", "uq(user_id, provider, alias)", ""],
        ["MariaDB", "llm_call_log", "idx(project_id), idx(purpose)", ""],
        ["MariaDB", "proposal_category", "uq(code)", ""],
        ["MariaDB", "audit_log", "idx(user_id), idx(action)", ""],
        ["MongoDB", "documents", "{projectUuid:1}, {attachmentId:1}", ""],
        ["MongoDB", "analysisResults", "{projectUuid:1}", ""],
        ["MongoDB", "llmSessions", "{projectUuid:1, createdAt:-1}, {purpose:1}", "TTL 90일"],
        ["MongoDB", "proposalDrafts", "{projectUuid:1, version:-1}", ""],
        ["MongoDB", "wbsTasks", "{projectUuid:1, version:-1}", ""],
        ["MongoDB", "categoryPrompts", "{code:1, version:-1}, partial(active=true)", ""],
    ], widths=[12, 22, 60, 18])

    # 시드
    ws = wb.create_sheet("21_시드데이터")
    r = add_title(ws, "표준 목차 시드 데이터", "proposal_category")
    write_table(ws, r, ["code", "name_ko", "sort_order", "is_active"], [
        ["OVERVIEW", "사업 개요", 10, 1],
        ["GENERAL", "일반 사항", 20, 1],
        ["TECH_REQ", "기술 요구사항", 30, 1],
        ["PM_REQ", "사업관리 요구사항", 40, 1],
        ["SECURITY", "보안 요구사항", 50, 1],
        ["CONSTRAINT", "제약 조건", 60, 1],
        ["ETC", "기타", 90, 1],
    ], widths=[18, 20, 12, 12])

    out = DOCS_DIR / "db-design.xlsx"
    wb.save(out)
    print(f"OK  -> {out.name}  ({len(wb.sheetnames)} sheets)")


# ╔══════════════════════════════════════════════════════════
# ║ 3) interface-spec.xlsx
# ╚══════════════════════════════════════════════════════════
def gen_iface_xlsx():
    wb = Workbook()
    ws = wb.active
    ws.title = "00_요약"
    r = add_title(ws, "Lon — 연계 정의서", "외부 LLM 3사 + 내부 REST + 시스템/DB 연계")
    write_table(ws, r, ["구분", "내용"], [
        ["통신", "TLS 1.2+, 외부 호출 60s 타임아웃, 재시도 3회 (1·2·4s)"],
        ["포맷", "JSON UTF-8, 시간 ISO-8601 UTC"],
        ["오류", "BE 코드 LON-XXXX + HTTP 상태 동시 반환"],
        ["로깅", "모든 외부 호출은 llm_call_log + llmSessions"],
        ["보안", "PII 마스킹(휴대폰/주민/이메일), 외부 도메인 화이트리스트"],
    ], widths=[12, 100])

    # 연계 일람
    ws = wb.create_sheet("01_연계일람")
    r = add_title(ws, "연계 대상 일람", "9건")
    write_table(ws, r, ["ID", "구분", "시스템", "방향", "프로토콜", "인증", "동기/비동기"], [
        ["IF-001", "외부", "OpenAI ChatGPT API", "OUT", "HTTPS/REST", "Bearer (API Key)", "동기"],
        ["IF-002", "외부", "Google Gemini API", "OUT", "HTTPS/REST", "x-goog-api-key", "동기"],
        ["IF-003", "외부", "Anthropic Claude API", "OUT", "HTTPS/REST", "x-api-key", "동기"],
        ["IF-101", "내부", "FE ↔ BE", "IN/OUT", "HTTPS/REST(JSON)", "Cookie/Session JWT", "동기"],
        ["IF-201", "시스템", "OS 키체인 (Windows Credential Manager)", "IN/OUT", "OS API", "OS 사용자", "동기"],
        ["IF-202", "시스템", "로컬 파일 시스템", "IN/OUT", "FS", "OS 권한", "동기"],
        ["IF-301", "데이터", "MariaDB 11.4", "IN/OUT", "TCP/3306", "user/password", "동기"],
        ["IF-302", "데이터", "MongoDB 8.2", "IN/OUT", "TCP/27017", "SCRAM", "동기"],
        ["IF-401", "출력", "PPTX/XLSX 산출물", "OUT", "FS", "-", "비동기 (잡)"],
    ], widths=[10, 8, 38, 10, 22, 22, 12])

    # IF-001 OpenAI
    ws = wb.create_sheet("IF-001_OpenAI")
    r = add_title(ws, "IF-001 · OpenAI ChatGPT", "GEN_PPTX / GEN_WBS / ANALYZE / TEST")
    r = write_table(ws, r, ["항목", "값"], [
        ["Endpoint", "POST https://api.openai.com/v1/chat/completions"],
        ["인증", "Authorization: Bearer {API_KEY}"],
        ["Content-Type", "application/json"],
        ["타임아웃", "60s"],
        ["재시도", "5xx/429 → 3회 (1s, 2s, 4s)"],
        ["모델", "gpt-4o, gpt-4o-mini"],
    ], widths=[16, 90])
    r += 1
    ws.cell(row=r, column=1, value="요청 예 (JSON)").font = Font(name=KO_FONT, bold=True, color=NAVY, size=11)
    r += 1
    req = '{"model":"gpt-4o-mini","messages":[{"role":"system","content":"<system>"},{"role":"user","content":"<프롬프트>"}],"temperature":0.4,"max_tokens":4096}'
    ws.cell(row=r, column=1, value=req).font = Font(name="Consolas", size=10)
    ws.cell(row=r, column=1).alignment = Alignment(wrap_text=True)
    r += 2
    write_table(ws, r, ["응답 키", "저장 위치"], [
        ["choices[0].message.content", "llmSessions.response.text"],
        ["usage.prompt_tokens", "llm_call_log.input_tokens"],
        ["usage.completion_tokens", "llm_call_log.output_tokens"],
    ], widths=[40, 40])

    # IF-002 Gemini
    ws = wb.create_sheet("IF-002_Gemini")
    r = add_title(ws, "IF-002 · Google Gemini", "기본 Provider")
    r = write_table(ws, r, ["항목", "값"], [
        ["Endpoint", "POST https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"],
        ["인증", "Header x-goog-api-key: {API_KEY} (or query ?key=)"],
        ["Content-Type", "application/json"],
        ["타임아웃", "60s"],
        ["모델", "gemini-2.5-flash, gemini-2.5-pro"],
    ], widths=[16, 100])
    r += 1
    ws.cell(row=r, column=1, value="요청 예").font = Font(name=KO_FONT, bold=True, color=NAVY, size=11)
    r += 1
    ws.cell(row=r, column=1, value='{"contents":[{"parts":[{"text":"<프롬프트>"}]}],"generationConfig":{"temperature":0.4,"maxOutputTokens":4096}}').font = Font(name="Consolas", size=10)

    # IF-003 Anthropic
    ws = wb.create_sheet("IF-003_Anthropic")
    r = add_title(ws, "IF-003 · Anthropic Claude", "")
    r = write_table(ws, r, ["항목", "값"], [
        ["Endpoint", "POST https://api.anthropic.com/v1/messages"],
        ["인증", "x-api-key: {API_KEY}, anthropic-version: 2023-06-01"],
        ["모델", "claude-opus-4-7, claude-sonnet-4-6"],
        ["타임아웃", "60s"],
    ], widths=[16, 100])
    r += 1
    ws.cell(row=r, column=1, value="요청 예").font = Font(name=KO_FONT, bold=True, color=NAVY, size=11)
    r += 1
    ws.cell(row=r, column=1, value='{"model":"claude-sonnet-4-6","max_tokens":4096,"system":"<system>","messages":[{"role":"user","content":"<프롬프트>"}]}').font = Font(name="Consolas", size=10)

    # 오류 코드 매핑
    ws = wb.create_sheet("04_오류코드")
    r = add_title(ws, "외부→내부 오류 코드 매핑", "")
    write_table(ws, r, ["외부", "의미", "내부 코드", "사용자 메시지"], [
        ["401", "인증 실패", "LON-LLM-401", "API 키가 유효하지 않습니다"],
        ["403", "권한/지역 차단", "LON-LLM-403", "사용 권한이 없습니다"],
        ["429", "Rate Limit", "LON-LLM-429", "잠시 후 다시 시도해 주세요"],
        ["5xx", "서버 오류", "LON-LLM-5XX", "LLM 서비스 일시 장애"],
        ["Timeout", "타임아웃", "LON-LLM-TO", "응답 지연. 재시도 권장"],
    ], widths=[12, 18, 16, 60])

    # IF-101 내부 API 목록
    ws = wb.create_sheet("IF-101_내부API")
    r = add_title(ws, "IF-101 · 내부 REST API", "FE ↔ BE")
    write_table(ws, r, ["#", "메서드", "경로", "요청", "응답", "권한", "부수효과"], [
        [1, "POST", "/api/llm/test", "{provider, model, apiKey}", "{ok, latencyMs, echo}", "USER", "-"],
        [2, "POST", "/api/files/analyze", "multipart: notice, references[]", "{projectUuid, documents, fields, confidence}", "USER", "documents+analysisResults+llm_call_log INSERT"],
        [3, "POST", "/api/projects", "Project JSON", "{id, uuid}", "USER", "project INSERT"],
        [4, "GET", "/api/projects/:uuid", "-", "Project + 산출물/로그", "USER", "-"],
        [5, "POST", "/api/generate/pptx", "{projectUuid, categories[]}", "PPTX 파일", "USER", "proposalDrafts + artifact INSERT"],
        [6, "POST", "/api/generate/wbs", "{projectUuid, phases}", "XLSX 파일", "USER", "wbsTasks + artifact INSERT"],
        [7, "GET", "/api/categories", "-", "표준 목차 트리", "USER", "-"],
    ], widths=[5, 8, 26, 36, 38, 8, 38])

    # 시스템 연계
    ws = wb.create_sheet("05_시스템연계")
    r = add_title(ws, "시스템 연계", "OS 키체인 / 파일 시스템")
    write_table(ws, r, ["ID", "대상", "용도", "라이브러리", "정책"], [
        ["IF-201", "Windows Credential Manager", "API Key 1차 저장소", "keyring (Python)", "키 네이밍 Lon/{provider}/{userUuid}"],
        ["IF-201F", "MariaDB ai_provider_setting", "키체인 사용 불가시 폴백", "cryptography(AESGCM)", "AES-256-GCM"],
        ["IF-202", "로컬 파일 시스템", "첨부/산출물 보관", "pathlib", "기본 폴더 %USERPROFILE%/Lon/workspace, sanitize 필수, 디스크<200MB 시 차단"],
    ], widths=[10, 30, 24, 22, 60])

    # DB 연계
    ws = wb.create_sheet("06_DB연계")
    r = add_title(ws, "데이터 저장소 연계", "MariaDB / MongoDB / Hybrid 패턴")
    r = write_table(ws, r, ["ID", "DB", "연결", "드라이버", "정책"], [
        ["IF-301", "MariaDB 11.4", "localhost:3306, lon", "mariadb / pymysql", "pool=10, recycle=1800s, READ COMMITTED"],
        ["IF-302", "MongoDB 8.2", "mongodb://lon_app@localhost:27017/lon", "pymongo>=4.7", "w=majority, j=true"],
    ], widths=[10, 14, 40, 24, 50])
    r += 1
    write_table(ws, r, ["단계", "동작"], [
        [1, "MariaDB INSERT ... status='PENDING'"],
        [2, "Mongo insertOne(...)"],
        [3, "MariaDB UPDATE ... status='OK', mongo_*_id=..."],
        ["실패", "status='ERROR' → mongo_repair 5분 주기 보상 잡"],
    ], widths=[8, 90])

    # 잡(비동기)
    ws = wb.create_sheet("07_비동기잡")
    r = add_title(ws, "비동기 잡 (IF-401 외)", "apscheduler + sqlite job store")
    write_table(ws, r, ["잡", "트리거", "처리"], [
        ["pptx-generate", "/api/generate/pptx", "python-pptx로 작성 → outputs/ 저장 → 완료 알림"],
        ["xlsx-generate", "/api/generate/wbs", "openpyxl로 WBS 생성"],
        ["mongo-repair", "5분 주기", "status='ERROR' 행 보상"],
        ["attachment-cleanup", "1시간 주기", "24h 경과 첨부 삭제"],
    ], widths=[22, 22, 70])

    # 모니터링
    ws = wb.create_sheet("08_모니터링")
    r = add_title(ws, "모니터링 지표 / 임계", "")
    write_table(ws, r, ["지표", "임계", "수집"], [
        ["LLM 평균 지연", "< 8s", "llm_call_log.latency_ms"],
        ["LLM 오류율 (일)", "< 2%", "http_status, error_code 집계"],
        ["DB 풀 사용률", "< 70%", "BE 메트릭"],
        ["디스크 사용률", "< 80%", "OS"],
        ["잡 실패율", "< 1%", "apscheduler 로그"],
    ], widths=[24, 18, 40])

    # 환경별 엔드포인트
    ws = wb.create_sheet("09_환경별엔드포인트")
    r = add_title(ws, "환경별 엔드포인트", "")
    write_table(ws, r, ["환경", "FE", "BE", "MariaDB", "MongoDB"], [
        ["로컬", "http://localhost:5173", "http://localhost:8080", "localhost:3306", "localhost:27017"],
        ["사내", "https://lon.intra", "https://lon-api.intra", "mariadb.intra:3306", "mongo.intra:27017"],
    ], widths=[10, 28, 28, 24, 24])

    out = DOCS_DIR / "interface-spec.xlsx"
    wb.save(out)
    print(f"OK  -> {out.name}  ({len(wb.sheetnames)} sheets)")


# ╔══════════════════════════════════════════════════════════
# ║ 4) api-architecture.xlsx
# ╚══════════════════════════════════════════════════════════
def gen_api_xlsx():
    wb = Workbook()

    # 표지
    ws = wb.active
    ws.title = "00_표지"
    r = add_title(ws, "Lon — API 구성도 (API Architecture)", "도메인/엔드포인트/의존성 매트릭스")
    write_table(ws, r, ["항목", "값"], [
        ["기준 BE", "FastAPI (Python 3.12)"],
        ["베이스 URL", "http://localhost:8080/api"],
        ["인증", "Cookie/Session JWT (HttpOnly + Secure)"],
        ["응답 포맷", "{ data, error, traceId }"],
        ["문서", "/docs (OpenAPI), /redoc"],
        ["도메인", "auth, projects, files, llm, generate, categories, admin"],
    ], widths=[14, 80])

    # 도메인 그룹
    ws = wb.create_sheet("01_도메인")
    r = add_title(ws, "API 도메인", "7 도메인")
    write_table(ws, r, ["도메인", "Prefix", "책임", "주요 엔드포인트 수"], [
        ["auth", "/api/auth", "로그인/세션/토큰", 3],
        ["projects", "/api/projects", "사업 메타 CRUD", 5],
        ["files", "/api/files", "첨부 업로드/분석", 2],
        ["llm", "/api/llm", "LLM 키 검증/상태", 2],
        ["generate", "/api/generate", "PPTX/XLSX 생성", 2],
        ["categories", "/api/categories", "표준 목차 마스터", 1],
        ["admin", "/api/admin", "사용자/프롬프트/감사 로그", 5],
    ], widths=[14, 20, 36, 14])

    # 전체 엔드포인트
    ws = wb.create_sheet("02_엔드포인트")
    r = add_title(ws, "전체 엔드포인트 카탈로그", "v0.1")
    rows = [
        # auth
        ["AUTH-01", "auth", "POST", "/api/auth/login", "로그인", "USER", "Y", "-", "-", "Maria", "-", "-", "-"],
        ["AUTH-02", "auth", "POST", "/api/auth/logout", "로그아웃", "USER", "Y", "-", "-", "-", "-", "-", "-"],
        ["AUTH-03", "auth", "POST", "/api/auth/refresh", "토큰 갱신", "USER", "Y", "-", "-", "Maria", "-", "-", "-"],
        # projects
        ["PROJ-01", "projects", "GET", "/api/projects", "사업 목록 조회", "USER", "Y", "-", "-", "Maria", "-", "-", "-"],
        ["PROJ-02", "projects", "POST", "/api/projects", "사업 생성", "USER", "Y", "-", "-", "Maria", "-", "audit_log", "-"],
        ["PROJ-03", "projects", "GET", "/api/projects/{uuid}", "사업 상세", "USER", "Y", "-", "-", "Maria", "Mongo", "-", "-"],
        ["PROJ-04", "projects", "PUT", "/api/projects/{uuid}", "사업 수정", "USER", "Y", "-", "-", "Maria", "-", "audit_log", "-"],
        ["PROJ-05", "projects", "DELETE", "/api/projects/{uuid}", "사업 삭제(논리)", "USER", "Y", "-", "-", "Maria", "-", "audit_log", "-"],
        # files
        ["FILE-01", "files", "POST", "/api/files/analyze", "공고문/산출물 분석", "USER", "Y", "Y(LLM)", "≤30MB", "Maria", "Mongo", "llm_call_log", "OUT→LLM"],
        ["FILE-02", "files", "DELETE", "/api/files/{id}", "첨부 삭제", "USER", "Y", "-", "-", "Maria", "Mongo", "-", "FS"],
        # llm
        ["LLM-01", "llm", "POST", "/api/llm/test", "API 키 검증", "USER", "Y", "Y(1tkn)", "-", "-", "-", "llm_call_log", "OUT→LLM"],
        ["LLM-02", "llm", "GET", "/api/llm/status", "사용량/상태", "USER", "Y", "-", "-", "Maria", "Mongo", "-", "-"],
        # generate
        ["GEN-01", "generate", "POST", "/api/generate/pptx", "사업제안서 PPTX 생성", "USER", "Y", "Y(LLM)", "비동기잡", "Maria", "Mongo", "artifact", "OUT→LLM, FS"],
        ["GEN-02", "generate", "POST", "/api/generate/wbs", "WBS XLSX 생성", "USER", "Y", "Y(LLM)", "비동기잡", "Maria", "Mongo", "artifact", "OUT→LLM, FS"],
        # categories
        ["CAT-01", "categories", "GET", "/api/categories", "표준 목차 트리", "USER", "Y", "-", "캐시", "Maria", "-", "-", "-"],
        # admin
        ["ADMIN-01", "admin", "GET", "/api/admin/users", "사용자 목록", "ADMIN", "Y", "-", "-", "Maria", "-", "-", "-"],
        ["ADMIN-02", "admin", "PUT", "/api/admin/users/{id}", "권한 변경", "ADMIN", "Y", "-", "-", "Maria", "-", "audit_log", "-"],
        ["ADMIN-03", "admin", "GET", "/api/admin/audit", "감사 로그 조회", "ADMIN", "Y", "-", "기간 검색", "Maria", "-", "-", "-"],
        ["ADMIN-04", "admin", "GET", "/api/admin/jobs", "잡 모니터", "ADMIN", "Y", "-", "-", "-", "-", "-", "scheduler"],
        ["ADMIN-05", "admin", "PUT", "/api/admin/category/{code}", "카테고리 수정", "ADMIN", "Y", "-", "-", "Maria", "Mongo", "audit_log", "-"],
    ]
    write_table(ws, r,
        ["ID","도메인","Method","Path","설명","권한","JWT","LLM 사용","비고","MariaDB","MongoDB","감사","외부"],
        rows,
        widths=[10,12,8,32,28,8,6,10,16,10,10,12,12]
    )

    # 권한 매트릭스
    ws = wb.create_sheet("03_권한매트릭스")
    r = add_title(ws, "권한 매트릭스", "도메인 × 역할")
    write_table(ws, r,
        ["도메인", "GUEST", "USER", "ADMIN"],
        [
            ["auth", "POST /login", "POST /logout, /refresh", "동일"],
            ["projects", "-", "ALL (자기 데이터)", "ALL (전체)"],
            ["files", "-", "ALL (자기 데이터)", "동일"],
            ["llm", "-", "test, status", "동일"],
            ["generate", "-", "pptx, wbs", "동일"],
            ["categories", "-", "GET", "GET / PUT (admin route)"],
            ["admin", "-", "-", "users, audit, jobs, category"],
        ],
        widths=[14, 14, 36, 36],
    )

    # 의존성 매트릭스
    ws = wb.create_sheet("04_의존성")
    r = add_title(ws, "엔드포인트 ↔ 자원 의존성", "MariaDB / MongoDB / FS / LLM")
    write_table(ws, r,
        ["ID", "MariaDB Tables", "Mongo Collections", "FS", "외부 LLM"],
        [
            ["AUTH-01", "user", "-", "-", "-"],
            ["PROJ-02", "project, audit_log", "-", "-", "-"],
            ["FILE-01", "project_attachment, llm_call_log", "documents, analysisResults, llmSessions", "attachments/", "Y (1회)"],
            ["LLM-01", "llm_call_log", "llmSessions", "-", "Y (1tkn)"],
            ["GEN-01", "artifact, llm_call_log", "proposalDrafts, llmSessions, categoryPrompts", "outputs/<uuid>/pptx/v{n}.pptx", "Y"],
            ["GEN-02", "artifact, llm_call_log", "wbsTasks, llmSessions", "outputs/<uuid>/xlsx/v{n}.xlsx", "Y"],
            ["CAT-01", "proposal_category", "-", "-", "-"],
            ["ADMIN-03", "audit_log, user", "-", "-", "-"],
        ],
        widths=[10, 30, 40, 30, 12],
    )

    # 호출 시퀀스: PPTX 생성
    ws = wb.create_sheet("05_시퀀스_PPTX")
    r = add_title(ws, "호출 시퀀스 — 사업제안서 PPTX 생성", "GEN-01")
    write_table(ws, r,
        ["#", "Actor", "Action", "대상", "비고"],
        [
            [1, "FE", "POST /api/generate/pptx { projectUuid, categories[] }", "BE", "체크박스 선택"],
            [2, "BE", "SELECT project, ai_provider_setting", "MariaDB", ""],
            [3, "BE", "load category prompts", "Mongo categoryPrompts", "active 버전"],
            [4, "BE", "loop categories → call LLM", "OpenAI/Gemini/Anthropic", "병렬 옵션"],
            [5, "BE", "INSERT llm_call_log + llmSessions", "Maria + Mongo", "Outbox 패턴"],
            [6, "BE", "INSERT proposalDrafts (categories.slides)", "Mongo", ""],
            [7, "BE", "render python-pptx", "FS outputs/<uuid>/pptx/v{n}.pptx", ""],
            [8, "BE", "INSERT artifact (type=PPTX, version=n)", "MariaDB", ""],
            [9, "BE", "응답 (200 file or {url})", "FE", "다운로드"],
            [10, "FE", "S-120 미리보기/다운로드", "User", ""],
        ],
        widths=[5, 8, 60, 30, 18],
    )

    # 호출 시퀀스: 분석
    ws = wb.create_sheet("06_시퀀스_분석")
    r = add_title(ws, "호출 시퀀스 — 첨부 분석", "FILE-01")
    write_table(ws, r,
        ["#", "Actor", "Action", "대상", "비고"],
        [
            [1, "FE", "POST /api/files/analyze (multipart)", "BE", ""],
            [2, "BE", "검증(ext, size, sha256)", "내부", "중복 차단"],
            [3, "BE", "save FS attachments/<uuid>/", "FS", ""],
            [4, "BE", "INSERT project_attachment(status=PENDING)", "MariaDB", ""],
            [5, "BE", "텍스트 추출(pdfminer/python-docx)", "내부", ""],
            [6, "BE", "INSERT documents (extractedText, chunks)", "Mongo", ""],
            [7, "BE", "LLM 호출 (요약 + 12필드 추출)", "Gemini 등", ""],
            [8, "BE", "INSERT llmSessions + analysisResults", "Mongo", ""],
            [9, "BE", "UPDATE project_attachment(status=OK, mongo_doc_id)", "MariaDB", ""],
            [10, "BE", "응답 (fields, confidence)", "FE", "③ 폼 자동 채움"],
        ],
        widths=[5, 8, 60, 30, 18],
    )

    # OpenAPI 컴포넌트(요약)
    ws = wb.create_sheet("07_스키마요약")
    r = add_title(ws, "OpenAPI 주요 스키마", "")
    write_table(ws, r,
        ["스키마", "필드", "타입", "비고"],
        [
            ["Project", "uuid", "string(uuid)", "PK"],
            ["Project", "companyName", "string", ""],
            ["Project", "projectName", "string", "필수"],
            ["Project", "goal/scope/...", "string", "12 필드"],
            ["Project", "aiProvider", "enum", "OPENAI/GEMINI/ANTHROPIC"],
            ["Project", "aiModel", "string", ""],
            ["Project", "status", "enum", "DRAFT/READY/GENERATED/ARCHIVED"],
            ["AnalyzeResponse", "projectUuid", "string", ""],
            ["AnalyzeResponse", "documents[]", "array<{id, slot}>", ""],
            ["AnalyzeResponse", "fields", "object", "12 필드 자동 채움"],
            ["AnalyzeResponse", "confidence", "object", "필드별 신뢰도"],
            ["GenerateRequest", "projectUuid", "string", ""],
            ["GenerateRequest", "categories[]", "array<string>", "code"],
            ["ErrorResponse", "code", "string", "LON-..."],
            ["ErrorResponse", "message", "string", ""],
            ["ErrorResponse", "details", "object", ""],
            ["ErrorResponse", "traceId", "string(uuid)", ""],
        ],
        widths=[22, 22, 22, 36],
    )

    # 비기능 (Rate Limit / 캐싱)
    ws = wb.create_sheet("08_비기능")
    r = add_title(ws, "API 비기능 정책", "")
    write_table(ws, r,
        ["항목", "정책"],
        [
            ["Rate Limit", "사용자당 60 req/min, 생성 API 6 req/min"],
            ["응답 시간 SLA", "조회<300ms / 분석<10s / 생성<15s (요건 30개)"],
            ["캐싱", "GET /categories → 10분 캐시"],
            ["페이지네이션", "?page=&size=&sort=updatedAt,desc"],
            ["멱등성", "POST /generate → Idempotency-Key 헤더"],
            ["관측", "X-Trace-Id 모든 응답 헤더에 포함"],
            ["보안", "CORS allowlist, CSRF 더블 서브밋"],
        ],
        widths=[20, 90],
    )

    out = DOCS_DIR / "api-architecture.xlsx"
    wb.save(out)
    print(f"OK  -> {out.name}  ({len(wb.sheetnames)} sheets)")


# ── main ──────────────────────────────────────────────────
if __name__ == "__main__":
    gen_screen_pptx()
    gen_db_xlsx()
    gen_iface_xlsx()
    gen_api_xlsx()
    print("DONE")
