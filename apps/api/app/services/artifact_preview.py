"""산출물 미리보기 — PPTX 슬라이드 / XLSX 시트 파싱."""
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation


def preview_pptx(path: Path, max_slides: int = 30) -> dict:
    prs = Presentation(str(path))
    slides_data = []
    total = 0
    for i, slide in enumerate(prs.slides):
        total += 1
        if i >= max_slides:
            continue
        title = ""
        bullets: list[str] = []
        # 슬라이드의 텍스트 셰이프 순회 (제목 placeholder 우선)
        title_set = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # placeholder type 으로 제목 추정
            try:
                ph_type = shape.placeholder_format.type if shape.is_placeholder else None
            except Exception:  # noqa: BLE001
                ph_type = None
            for para in tf.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if not txt:
                    continue
                # placeholder type 1 = TITLE, 13 = CENTER_TITLE
                if not title_set and ph_type in (1, 13):
                    title = txt
                    title_set = True
                elif not title_set and not title:
                    title = txt
                    title_set = True
                else:
                    bullets.append(txt)
        note = ""
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:  # noqa: BLE001
            pass
        slides_data.append({
            "index": i + 1,
            "title": title[:200],
            "bullets": [b[:300] for b in bullets[:12]],
            "speakerNote": note[:600],
        })
    return {
        "format": "PPTX",
        "totalSlides": total,
        "shownSlides": len(slides_data),
        "slides": slides_data,
    }


def preview_xlsx(path: Path, max_rows: int = 50, max_cols: int = 20, max_sheets: int = 10) -> dict:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        sheets = []
        names = wb.sheetnames[:max_sheets]
        for name in names:
            ws = wb[name]
            rows: list[list[str]] = []
            r_count = 0
            for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if r_idx >= max_rows:
                    break
                cells = []
                for c in row[:max_cols]:
                    if c is None:
                        cells.append("")
                    elif isinstance(c, (int, float)):
                        cells.append(str(c))
                    else:
                        cells.append(str(c)[:200])
                rows.append(cells)
                r_count = r_idx + 1
            sheets.append({
                "name": name,
                "totalRows": ws.max_row or 0,
                "totalCols": ws.max_column or 0,
                "shownRows": r_count,
                "rows": rows,
            })
    finally:
        wb.close()
    return {
        "format": "XLSX",
        "totalSheets": len(wb.sheetnames),
        "shownSheets": len(sheets),
        "limit": {"rows": max_rows, "cols": max_cols, "sheets": max_sheets},
        "sheets": sheets,
    }
