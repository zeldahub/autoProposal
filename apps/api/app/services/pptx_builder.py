"""사업제안서 PPTX 빌더 (LLM-driven 또는 placeholder)."""
from typing import IO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GRAY = RGBColor(0x66, 0x66, 0x66)

DEFAULT_CATEGORIES = [
    ("OVERVIEW", "사업 개요"),
    ("GENERAL", "일반 사항"),
    ("TECH_REQ", "기술 요구사항"),
    ("PM_REQ", "사업관리 요구사항"),
    ("SECURITY", "보안 요구사항"),
    ("CONSTRAINT", "제약 조건"),
    ("ETC", "기타"),
]


def _add_title_slide(prs, project_name: str, company: str | None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(11.9), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "사업 제안서"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
    sub = s.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.7))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = project_name or "(사업명 미입력)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(22); p2.font.color.rgb = ACCENT
    if company:
        c = s.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.5))
        cp = c.text_frame.paragraphs[0]
        cp.text = company
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(14); cp.font.color.rgb = GRAY


def _add_section_divider(prs, name: str):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(11.9), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = NAVY


def _add_content_slide(prs, title: str, bullets: list[str], note: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 헤더
    head = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    p = head.text_frame.paragraphs[0]
    p.text = title or "(제목 없음)"
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = NAVY
    # accent bar
    # 본문
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets or []:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.text = f"• {b}"
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_after = Pt(8)
    # 발표자 노트
    if note:
        s.notes_slide.notes_text_frame.text = note


def build_proposal_pptx(
    out: IO[bytes],
    project_uuid: str,
    project: dict | None = None,
    drafts: list[dict] | None = None,
    categories: list[str] | None = None,
) -> None:
    """drafts(LLM 생성 결과)이 있으면 그것을, 없으면 placeholder 생성."""
    project = project or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, project.get("projectName") or f"Project {project_uuid[:8]}",
                     project.get("companyName"))

    if drafts:
        for d in drafts:
            name = d.get("name") or d.get("code", "")
            _add_section_divider(prs, name)
            for sl in d.get("slides", []) or []:
                _add_content_slide(
                    prs,
                    sl.get("title", name),
                    sl.get("bullets", []),
                    sl.get("speakerNote", ""),
                )
    else:
        cats = categories or [c[0] for c in DEFAULT_CATEGORIES]
        name_map = dict(DEFAULT_CATEGORIES)
        for code in cats:
            name = name_map.get(code, code)
            _add_section_divider(prs, name)
            _add_content_slide(prs, name, ["(LLM 미사용 — 수동 보완 필요)"], "")

    prs.save(out)
