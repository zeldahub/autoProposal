"""산출물 인라인 편집 — preview 와 동일한 트리(슬라이드#/시트/행/열) 좌표로 텍스트만 갱신.

다음 신규 버전을 별도 파일로 생성하여 원본 보존.
"""
from pathlib import Path
from typing import TypedDict

from openpyxl import load_workbook
from pptx import Presentation


class PptxBulletEdit(TypedDict, total=False):
    index: int          # 1-based slide index
    title: str | None   # None 이면 변경 안 함
    bullets: list[str] | None
    speakerNote: str | None


class XlsxCellEdit(TypedDict):
    sheet: str
    row: int   # 1-based
    col: int   # 1-based
    value: str


def apply_pptx_edits(src: Path, dst: Path, edits: list[PptxBulletEdit]) -> dict:
    """슬라이드별로 title/bullets/speakerNote 를 일괄 갱신.

    동작 규칙:
    - title 은 placeholder TITLE/CENTER_TITLE 셰이프의 첫 텍스트를 교체
    - bullets 는 그 외 텍스트 셰이프 중 첫 번째 (제목 외 본문 셰이프) 의 단락을 재구성
    - speakerNote 는 슬라이드 노트 텍스트 frame 전체를 교체
    """
    prs = Presentation(str(src))
    by_index = {e["index"]: e for e in edits if "index" in e}
    applied = {"slides": 0, "titles": 0, "bulletGroups": 0, "notes": 0}

    for i, slide in enumerate(prs.slides, start=1):
        edit = by_index.get(i)
        if not edit:
            continue
        applied["slides"] += 1

        title_shape = None
        body_shape = None

        # 1차: placeholder TITLE 우선
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                ph_type = shape.placeholder_format.type if shape.is_placeholder else None
            except Exception:  # noqa: BLE001
                ph_type = None
            if ph_type in (1, 13) and title_shape is None:
                title_shape = shape
                break

        # 2차: title_shape 가 없으면 첫 textbox 를 title, 다음 textbox 를 body 로
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape is title_shape:
                continue
            if title_shape is None:
                title_shape = shape
                continue
            if body_shape is None:
                body_shape = shape
                break

        new_title = edit.get("title")
        if new_title is not None and title_shape is not None:
            tf = title_shape.text_frame
            # 단순화: 모든 단락의 첫 번째 run 만 보존하고 텍스트 교체
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.runs:
                    p.runs[0].text = new_title
                    for r in p.runs[1:]:
                        r.text = ""
                else:
                    p.text = new_title
                # 두번째 이후 단락 삭제
                for p2 in list(tf.paragraphs[1:]):
                    p2._p.getparent().remove(p2._p)
            applied["titles"] += 1

        new_bullets = edit.get("bullets")
        if new_bullets is not None and body_shape is not None:
            tf = body_shape.text_frame
            tf.clear()
            # tf.clear() 후 paragraph 가 1개 남음 — 거기에 첫 줄 채우고 이후는 add_paragraph
            first = True
            for line in new_bullets:
                if first:
                    tf.paragraphs[0].text = line
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.text = line
            applied["bulletGroups"] += 1

        new_note = edit.get("speakerNote")
        if new_note is not None:
            try:
                ns = slide.notes_slide
                ns.notes_text_frame.text = new_note
                applied["notes"] += 1
            except Exception:  # noqa: BLE001
                pass

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    return applied


def apply_xlsx_edits(src: Path, dst: Path, edits: list[XlsxCellEdit]) -> dict:
    """시트/행/열 좌표의 셀 값을 문자열로 갱신."""
    wb = load_workbook(str(src))  # read_only=False (편집)
    try:
        applied = {"cells": 0, "missingSheets": []}
        for e in edits:
            name = e["sheet"]
            if name not in wb.sheetnames:
                if name not in applied["missingSheets"]:
                    applied["missingSheets"].append(name)
                continue
            ws = wb[name]
            row = max(1, int(e["row"]))
            col = max(1, int(e["col"]))
            ws.cell(row=row, column=col).value = e.get("value", "")
            applied["cells"] += 1

        dst.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(dst))
        return applied
    finally:
        wb.close()
