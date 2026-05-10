"""artifact_editor — PPTX/XLSX 인라인 편집 단위 테스트 (실제 파일 round-trip)."""
import io
from pathlib import Path

import pytest
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.artifact_editor import apply_pptx_edits, apply_xlsx_edits


def _make_pptx(path: Path) -> None:
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "원본 제목"
    body_box = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = body_box.text_frame
    tf.paragraphs[0].text = "원본 bullet 1"
    p = tf.add_paragraph()
    p.text = "원본 bullet 2"
    prs.save(str(path))


def _make_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "헤더1"
    ws["B1"] = "헤더2"
    ws["A2"] = "값1"
    ws["B2"] = "값2"
    wb2 = wb.create_sheet("Other")
    wb2["A1"] = "다른 시트"
    wb.save(str(path))
    wb.close()


class TestPptxEdit:
    def test_title_only(self, tmp_path: Path):
        src = tmp_path / "src.pptx"
        dst = tmp_path / "dst.pptx"
        _make_pptx(src)
        result = apply_pptx_edits(src, dst, [{"index": 1, "title": "새 제목"}])
        assert result["slides"] == 1
        assert result["titles"] == 1
        assert result["bulletGroups"] == 0

        prs = Presentation(str(dst))
        first_text = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                first_text.append(shape.text_frame.text)
        assert any("새 제목" in t for t in first_text)
        # 본문은 보존
        assert any("원본 bullet 1" in t for t in first_text)

    def test_bullets_only(self, tmp_path: Path):
        src = tmp_path / "src.pptx"
        dst = tmp_path / "dst.pptx"
        _make_pptx(src)
        result = apply_pptx_edits(src, dst, [{
            "index": 1, "bullets": ["새 bullet A", "새 bullet B", "새 bullet C"],
        }])
        assert result["bulletGroups"] == 1

        prs = Presentation(str(dst))
        full = "\n".join(s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame)
        assert "새 bullet A" in full
        assert "새 bullet B" in full
        assert "새 bullet C" in full
        # 원래 bullet 은 사라져야 함
        assert "원본 bullet 1" not in full

    def test_speaker_note(self, tmp_path: Path):
        src = tmp_path / "src.pptx"
        dst = tmp_path / "dst.pptx"
        _make_pptx(src)
        result = apply_pptx_edits(src, dst, [{"index": 1, "speakerNote": "노트 내용입니다"}])
        assert result["notes"] == 1
        prs = Presentation(str(dst))
        assert prs.slides[0].has_notes_slide
        assert "노트 내용" in prs.slides[0].notes_slide.notes_text_frame.text

    def test_missing_slide_index_skipped(self, tmp_path: Path):
        src = tmp_path / "src.pptx"
        dst = tmp_path / "dst.pptx"
        _make_pptx(src)
        # slide 99 는 없음 → 무시
        result = apply_pptx_edits(src, dst, [{"index": 99, "title": "X"}])
        assert result["slides"] == 0


class TestXlsxEdit:
    def test_basic_cell(self, tmp_path: Path):
        src = tmp_path / "src.xlsx"
        dst = tmp_path / "dst.xlsx"
        _make_xlsx(src)
        result = apply_xlsx_edits(src, dst, [
            {"sheet": "Data", "row": 1, "col": 1, "value": "수정된 헤더1"},
            {"sheet": "Data", "row": 2, "col": 2, "value": "수정된 값2"},
        ])
        assert result["cells"] == 2
        assert result["missingSheets"] == []

        wb = load_workbook(str(dst))
        ws = wb["Data"]
        assert ws["A1"].value == "수정된 헤더1"
        assert ws["B1"].value == "헤더2"  # 보존
        assert ws["B2"].value == "수정된 값2"
        # 다른 시트는 보존
        assert wb["Other"]["A1"].value == "다른 시트"
        wb.close()

    def test_missing_sheet_listed(self, tmp_path: Path):
        src = tmp_path / "src.xlsx"
        dst = tmp_path / "dst.xlsx"
        _make_xlsx(src)
        result = apply_xlsx_edits(src, dst, [
            {"sheet": "NoSuch", "row": 1, "col": 1, "value": "X"},
        ])
        assert result["cells"] == 0
        assert result["missingSheets"] == ["NoSuch"]
