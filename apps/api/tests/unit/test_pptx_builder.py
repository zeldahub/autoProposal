"""pptx_builder + xlsx_builder 단위 테스트 — 빌더가 유효한 파일을 만드는지 확인."""
import io

from openpyxl import load_workbook
from pptx import Presentation

from app.services.pptx_builder import build_proposal_pptx
from app.services.xlsx_builder import build_wbs_xlsx


class TestPptxBuilder:
    def test_placeholder_fallback(self):
        buf = io.BytesIO()
        build_proposal_pptx(
            buf,
            project_uuid="abc12345-aaaa-bbbb-cccc-1234567890ab",
            project={"projectName": "테스트 사업", "companyName": "ACME"},
        )
        buf.seek(0)
        prs = Presentation(buf)
        # title + 7 카테고리 × (divider+content) = 1 + 14 슬라이드
        assert len(prs.slides) >= 8

        # 첫 슬라이드에 사업명/회사명 포함
        first_text = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                first_text.append(shape.text_frame.text)
        joined = " ".join(first_text)
        assert "테스트 사업" in joined
        assert "ACME" in joined

    def test_with_drafts(self):
        drafts = [{
            "code": "OVERVIEW", "name": "사업 개요",
            "slides": [{"title": "S1", "bullets": ["a", "b"], "speakerNote": "n1"}],
        }]
        buf = io.BytesIO()
        build_proposal_pptx(
            buf, project_uuid="x", project={"projectName": "P"},
            drafts=drafts,
        )
        buf.seek(0)
        prs = Presentation(buf)
        # title + divider + 1 content = 3 슬라이드
        assert len(prs.slides) == 3
        # 노트 검증
        last = prs.slides[2]
        assert last.has_notes_slide
        assert "n1" in last.notes_slide.notes_text_frame.text


class TestXlsxBuilder:
    def test_basic(self):
        buf = io.BytesIO()
        build_wbs_xlsx(buf, project_uuid="proj", phases=4)
        buf.seek(0)
        wb = load_workbook(buf)
        # 시트 존재
        assert len(wb.sheetnames) >= 1
        ws = wb[wb.sheetnames[0]]
        # 헤더 행이 존재
        assert ws.max_row >= 1
        wb.close()
