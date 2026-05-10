"""artifact_preview — preview_pptx / preview_xlsx 단위 테스트."""
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.services.artifact_preview import preview_pptx, preview_xlsx


def _make_pptx_with_n_slides(path: Path, n: int) -> None:
    prs = Presentation()
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title.text_frame.paragraphs[0].text = f"슬라이드 {i+1} 제목"
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        body.text_frame.paragraphs[0].text = f"본문 {i+1}-A"
        p = body.text_frame.add_paragraph()
        p.text = f"본문 {i+1}-B"
    prs.save(str(path))


class TestPreviewPptx:
    def test_basic(self, tmp_path: Path):
        f = tmp_path / "p.pptx"
        _make_pptx_with_n_slides(f, 3)
        data = preview_pptx(f)
        assert data["format"] == "PPTX"
        assert data["totalSlides"] == 3
        assert data["shownSlides"] == 3
        assert len(data["slides"]) == 3
        s0 = data["slides"][0]
        assert s0["index"] == 1
        assert "슬라이드 1 제목" in s0["title"]
        assert any("본문 1-A" in b for b in s0["bullets"])
        assert any("본문 1-B" in b for b in s0["bullets"])

    def test_max_slides_limit(self, tmp_path: Path):
        f = tmp_path / "p.pptx"
        _make_pptx_with_n_slides(f, 8)
        data = preview_pptx(f, max_slides=3)
        assert data["totalSlides"] == 8
        assert data["shownSlides"] == 3
        assert len(data["slides"]) == 3


class TestPreviewXlsx:
    def test_basic(self, tmp_path: Path):
        f = tmp_path / "x.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["A", "B", "C"])
        ws.append([1, 2, 3])
        ws.append([4, 5, 6])
        wb.create_sheet("Empty")
        wb.save(str(f))
        wb.close()

        data = preview_xlsx(f)
        assert data["format"] == "XLSX"
        assert data["totalSheets"] == 2
        s0 = next(s for s in data["sheets"] if s["name"] == "Sheet1")
        assert s0["totalRows"] >= 3
        assert s0["totalCols"] >= 3
        assert s0["rows"][0] == ["A", "B", "C"]
        assert s0["rows"][1] == ["1", "2", "3"]

    def test_truncation(self, tmp_path: Path):
        f = tmp_path / "big.xlsx"
        wb = Workbook()
        ws = wb.active
        for i in range(100):
            ws.append([i, i + 1, i + 2])
        wb.save(str(f))
        wb.close()

        data = preview_xlsx(f, max_rows=10)
        sheet = data["sheets"][0]
        assert sheet["shownRows"] == 10
        assert sheet["totalRows"] == 100
        assert len(sheet["rows"]) == 10
